import os
import win32com.client as win32
from openpyxl import load_workbook
from pptx import Presentation
from datetime import datetime

def convert_pptx_to_pdf(pptx_path, pdf_path):
    powerpoint = win32.Dispatch("PowerPoint.Application")
    ppt = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
    
    ppt.SaveAs(os.path.abspath(pdf_path), FileFormat=32)  
    ppt.Close()

    powerpoint.Quit()

excel_file = 'data.xlsx'
wb = load_workbook(excel_file)
sheet = wb.active


pptx_file = 'template.pptx'
prs = Presentation(pptx_file)


if not os.path.exists("PDFs"):
    os.makedirs("PDFs")


for row_index, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
    tag_2_value = row[1]  
    
   
    row_presentation = Presentation(pptx_file)

    for slide in row_presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for i, column_name in enumerate(sheet[1], start=2):
                        if column_name.value and column_name.value in run.text:
                            column_value = row[i - 2]  
                            if isinstance(column_value, datetime):
                                formatted_date = column_value.strftime("%d-%m-%Y")  
                                run.text = run.text.replace(column_name.value, str(formatted_date))
                            else:
                                run.text = run.text.replace(column_name.value, str(column_value))

   
    output_pptx_file = f'{tag_2_value}.pptx'
    row_presentation.save(output_pptx_file)
    print(f'Updated presentation for row {row_index} saved as {output_pptx_file}')

    
    pdf_file = os.path.join("PDFs", f'{tag_2_value}.pdf')
    convert_pptx_to_pdf(output_pptx_file, pdf_file)
    print(f'Converted {output_pptx_file} to {pdf_file}')
    
    
    os.remove(output_pptx_file)
    print(f'Deleted {output_pptx_file}')

print('All presentations updated, saved as PDFs, and converted PPTX files deleted.')
